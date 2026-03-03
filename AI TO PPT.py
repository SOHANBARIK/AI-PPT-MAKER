import streamlit as st
import base64
import pptx
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Inches
import os
import requests
from dotenv import load_dotenv
from openai import OpenAI
import re  

# Load environment variables for local development
load_dotenv()

# Safely fetch keys from Streamlit Secrets (Cloud) OR Environment Variables (Local)
try:
    GROQ_API_KEY = st.secrets["GROQ_API_KEY"]
    PEXELS_API_KEY = st.secrets["PEXELS_API_KEY"]
except (FileNotFoundError, KeyError):
    GROQ_API_KEY = os.getenv("GROQ_API_KEY")
    PEXELS_API_KEY = os.getenv("PEXELS_API_KEY")

# Initialize OpenAI client pointing to Groq's API
client = OpenAI(
    api_key=GROQ_API_KEY, 
    base_url="https://api.groq.com/openai/v1"
)

# -------------------- Helper Functions --------------------
def clean_font_selection(selection, label):
    """Handles font family selection and 'Other' input"""
    if selection == "Other":
        custom_font = st.text_input(f"Enter custom font family for {label}:").strip()
        if not custom_font or custom_font in ["Calibri", "Arial", "Times New Roman", "Verdana", "Tahoma", "Other"]:
            st.warning(f"⚠️ Please enter a valid custom font for {label}. Defaulting to Calibri.")
            return "Calibri"
        return custom_font
    return selection

def clean_text_and_tags(text):
    """Aggressively cleans AI 'robot noise', raw tokens, and conversational filler."""
    if not text:
        return ""
    
    text = re.sub(r'<s>|</s>|\[/?\w+\]', '', text)
    text = text.replace('</li>', '\n').replace('</p>', '\n').replace('<br>', '\n')
    text = re.sub(r'<[^>]+>', '', text)
    text = text.replace('**', '').replace('__', '')
    
    lines = text.split('\n')
    cleaned_lines = []
    for line in lines:
        lower_line = line.lower().strip()
        if lower_line.startswith(("here is", "sure", "certainly", "below are", "i have generated")):
            continue
        cleaned_lines.append(line)
    
    text = "\n".join(cleaned_lines)
    text = re.sub(r'\n\s*\n', '\n', text)
    return text.strip()

# ------------------ AI FUNCTIONS ------------------ #
def generate_slide_titles(topic, num_slides=5):
    """Generate slide titles with strict instruction to avoid conversational filler."""
    system_prompt = "You are a data generator. Return ONLY the slide titles. No introductions. No numbering. No 'Here are the titles'."
    user_prompt = f"Generate exactly {num_slides} short, professional slide titles for a presentation on '{topic}'."

    try:
        response = client.chat.completions.create(
            model="llama3-8b-8192",  # Updated to a fast Groq model
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=200,
        )
        
        raw_content = clean_text_and_tags(response.choices[0].message.content)
        raw_titles = raw_content.split("\n")
        titles = []
        for t in raw_titles:
            clean_t = re.sub(r'^\d+\.?\s*', '', t.strip())
            if clean_t and len(clean_t) > 2: 
                titles.append(clean_t)
                
        return titles[:num_slides]
    except Exception as e:
        st.error(f"Groq API Error generating titles: {e}")
        return []

def generate_slide_content(slide_title, style="bullets"):
    """Generate content with strict 'No Conversational Filler' rules."""
    if style == "bullets":
        user_prompt = f"Write 3-5 concise bullet points for a slide titled '{slide_title}'. Return ONLY the bullet points. Do not say 'Here are the points'."
    else:
        user_prompt = f"Write a short paragraph (4 sentences) for a slide titled '{slide_title}'. Return ONLY the paragraph."

    try:
        response = client.chat.completions.create(
            model="llama3-8b-8192",  # Updated to a fast Groq model
            messages=[
                {"role": "system", "content": "You are a strict content generator. Output ONLY the requested content. No conversational filler text."},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=300,
        )
        return clean_text_and_tags(response.choices[0].message.content)
    except Exception as e:
        st.error(f"Groq API Error generating content: {e}")
        return ""

def fetch_pexels_image(query, save_path):
    """Fetch an image from Pexels API based on query."""
    headers = {"Authorization": PEXELS_API_KEY}
    params = {"query": query, "per_page": 1, "orientation": "landscape"}
    url = "https://api.pexels.com/v1/search"

    try:
        response = requests.get(url, headers=headers, params=params)
        data = response.json()
        if "photos" in data and len(data["photos"]) > 0:
            image_url = data["photos"][0]["src"]["large"]
            img_data = requests.get(image_url).content
            with open(save_path, "wb") as f:
                f.write(img_data)
            return save_path
        else:
            return None
    except Exception as e:
        return None

# ------------------ PPT CREATION ------------------ #
def create_presentation(topic, slide_titles, slide_contents, style,
                        first_page_title_size, global_title_size, content_font_size, font_title, font_content, add_images=False):
    """Create a PPT with generated titles, content, and optional images."""
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    title_shape.text = clean_text_and_tags(topic)

    text_frame = title_shape.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(first_page_title_size)
            run.font.bold = True
            run.font.name = font_title

    # Content Slides
    for idx, (slide_title, slide_content) in enumerate(zip(slide_titles, slide_contents), start=1):
        slide = prs.slides.add_slide(slide_layout)

        # Slide Title
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        for p in title_shape.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(global_title_size)
                run.font.name = font_title

        # Content Area
        text_frame = slide.shapes.placeholders[1].text_frame
        text_frame.clear()

        if style == "bullets":
            for line in slide_content.split("\n"):
                line = line.strip("-• \t")
                if line:
                    p = text_frame.add_paragraph()
                    p.text = line
                    p.level = 0
                    p.font.size = Pt(content_font_size)
                    p.font.name = font_content
        else:  
            p = text_frame.paragraphs[0]
            p.text = slide_content
            p.font.size = Pt(content_font_size)

        # Optional: Add Pexels image
        if add_images and PEXELS_API_KEY:
            os.makedirs("generated_ppt/images", exist_ok=True)
            image_path = f"generated_ppt/images/slide_{idx}.jpg"
            if fetch_pexels_image(slide_title, image_path):
                try:
                    slide.shapes.add_picture(image_path, Inches(5), Inches(2), Inches(3), Inches(3))
                except Exception:
                    pass

    # Save PPT
    os.makedirs("generated_ppt", exist_ok=True)
    ppt_filename = f"generated_ppt/{topic.replace(' ', '_')}_presentation.pptx"
    prs.save(ppt_filename)
    return ppt_filename

# ------------------ STREAMLIT APP ------------------ #
def get_ppt_download_link(ppt_filename):
    """Generate download link for the PPT."""
    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{os.path.basename(ppt_filename)}">Download Presentation</a>'

def main():
    st.title("📊 AI-Powered PPT Maker")
    st.subheader("Generate PowerPoint Presentations with Pexels Images")

    topic = st.text_input("Enter the topic for your presentation:")
    num_slides = st.slider("Number of slides", min_value=2, max_value=10, value=3, step=1)
    first_page_title_size = st.slider("Title font size (pt):", min_value=10, max_value=60, value=32, step=2)
    global_title_size = st.slider("Slide titles font size:", 20, 60, 32, 2)
    content_font_size = st.slider("Content font size (pt):", min_value=8, max_value=40, value=16, step=1)
    
    st.write("Font Settings")
    same_font = st.radio("Use the same font for Title & Content?", ["Yes", "No"], index=0)

    if same_font == "Yes":
        font = st.selectbox("Choose font family:", ["Calibri", "Arial", "Times New Roman", "Verdana", "Tahoma", "Other"])
        font_title = font_content = clean_font_selection(font, "Slides")
    else:
        font_title = st.selectbox("Font family for Titles:", ["Calibri", "Arial", "Times New Roman", "Verdana", "Tahoma", "Other"])
        font_title = clean_font_selection(font_title, "Titles")
        font_content = st.selectbox("Font family for Content:", ["Calibri", "Arial", "Times New Roman", "Verdana", "Tahoma", "Other"])
        font_content = clean_font_selection(font_content, "Content")

    style = st.radio("Content style:", ["bullets", "paragraph"])
    add_images = st.checkbox("Include Pexels images in slides")
    generate_button = st.button("🚀 Generate Presentation")

    if generate_button and topic:
        st.info("🔍 Generating slide titles...")
        slide_titles = generate_slide_titles(topic, num_slides)
        
        # Check if titles generated successfully
        if not slide_titles:
            st.error("Failed to generate slide titles. Check API Key or connection.")
            return
            
        st.write("✅ Slide titles:", slide_titles)

        st.info("📝 Generating slide contents...")
        slide_contents = [generate_slide_content(title, style) for title in slide_titles]

        st.info("📂 Creating presentation...")
        ppt_filename = create_presentation(topic, slide_titles, slide_contents, style,first_page_title_size,global_title_size, content_font_size, font_title, font_content, add_images)

        st.success("🎉 Presentation generated successfully!")
        st.markdown(get_ppt_download_link(ppt_filename), unsafe_allow_html=True)

    st.markdown("""
        <style>
            body {
                background-color: rgb(1, 1, 1);
                color: white;
                font-family: 'Arial', sans-serif;
            }
            .fixed-bottom {
                position: fixed;
                bottom: 10px;
                left: 50%;
                transform: translateX(-50%);
                font-size: 12px;
                color: white;
                background-color: rgb(0, 0, 0);
                padding: 5px 10px;
                border-radius: 5px;
                z-index: 1000;
            }    
        </style>
        <div class="fixed-bottom">made with ❤️ from Sohan </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
